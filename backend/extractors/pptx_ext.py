"""PPTX text extraction using python-pptx.

Walks every slide, every shape (including grouped), every text frame,
table cells, slide notes, and core properties. Returns slide count.
"""
from __future__ import annotations

from pathlib import Path

from . import ExtractionResult


def _iter_shapes(shapes):
    """Yield shapes recursively (descend into group shapes)."""
    for shape in shapes:
        yield shape
        try:
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                yield from _iter_shapes(shape.shapes)
        except Exception:
            pass


def _shape_text(shape) -> str:
    parts: list[str] = []
    # Text frames
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text or "" for run in para.runs).strip()
            if line:
                parts.append(line)
    # Tables
    if getattr(shape, "has_table", False):
        try:
            for row in shape.table.rows:
                cells = [(c.text or "").strip() for c in row.cells]
                row_text = " | ".join(c for c in cells if c)
                if row_text:
                    parts.append(row_text)
        except Exception:
            pass
    return "\n".join(parts)


def extract_pptx(path: Path) -> ExtractionResult:
    from pptx import Presentation

    prs = Presentation(str(path))
    slide_count = len(prs.slides)

    parts: list[str] = []

    # Core properties first
    cp = prs.core_properties
    if cp.title:
        parts.append(f"Title (metadata): {cp.title}")
    if cp.author:
        parts.append(f"Author (metadata): {cp.author}")
    if cp.last_modified_by:
        parts.append(f"Last modified by: {cp.last_modified_by}")
    if cp.subject:
        parts.append(f"Subject (metadata): {cp.subject}")
    if cp.keywords:
        parts.append(f"Keywords (metadata): {cp.keywords}")
    if cp.created:
        parts.append(f"Created (metadata): {cp.created}")
    if cp.modified:
        parts.append(f"Modified (metadata): {cp.modified}")
    if cp.revision:
        parts.append(f"Revision (metadata): {cp.revision}")

    for idx, slide in enumerate(prs.slides, 1):
        slide_lines: list[str] = [f"--- Slide {idx} ---"]
        for shape in _iter_shapes(slide.shapes):
            t = _shape_text(shape)
            if t:
                slide_lines.append(t)
        # Slide notes
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    slide_lines.append(f"[NOTES] {notes}")
        except Exception:
            pass

        if len(slide_lines) > 1:
            parts.append("\n".join(slide_lines))

    text = "\n\n".join(parts).strip()
    return ExtractionResult(text=text, page_count=slide_count)